from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

NAVY   = RGBColor(0x1A, 0x2F, 0x6E)
RED    = RGBColor(0xE8, 0x19, 0x2C)
BLUE   = RGBColor(0x1D, 0x4F, 0xA8)
ORANGE = RGBColor(0xE0, 0x70, 0x00)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GRY    = RGBColor(0xE2, 0xE2, 0xE2)
TALT   = RGBColor(0xEE, 0xF3, 0xFA)
TLINE  = RGBColor(0xC5, 0xD2, 0xE8)
FBKG   = RGBColor(0xF4, 0xF7, 0xFB)
DRK    = RGBColor(0x1A, 0x1A, 0x35)
MDM    = RGBColor(0x44, 0x50, 0x68)
LGT    = RGBColor(0x88, 0x95, 0xB0)
FONT   = "Calibri"

def C(r, g, b): return RGBColor(r, g, b)

def no_line(shape):
    spPr = shape._element.spPr
    for ln in spPr.findall(qn('a:ln')): spPr.remove(ln)
    spPr.append(parse_xml(
        '<a:ln xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        '<a:noFill/></a:ln>'))

def R(sl, x, y, w, h, fill=None, lc=None, lw=0.5):
    s = sl.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill: s.fill.solid(); s.fill.fore_color.rgb = fill
    else: s.fill.background()
    if lc: s.line.color.rgb = lc; s.line.width = Pt(lw)
    else: no_line(s)
    return s

def O(sl, x, y, w, h, lc=None, lw=2.5):
    s = sl.shapes.add_shape(9, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.background()
    if lc: s.line.color.rgb = lc; s.line.width = Pt(lw)
    else: no_line(s)
    return s

def T(sl, text, x, y, w, h, sz=11, bold=False, italic=False,
      color=DRK, align=PP_ALIGN.LEFT):
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run()
    r.text = text; r.font.name = FONT; r.font.size = Pt(sz)
    r.font.bold = bold; r.font.italic = italic; r.font.color.rgb = color
    return tb

def ML(sl, lines, x, y, w, h, dsz=11, dc=DRK):
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for ln in lines:
        if isinstance(ln, str): ln = {'t': ln}
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = ln.get('align', PP_ALIGN.LEFT)
        rn = p.add_run()
        rn.text = ln.get('t', ''); rn.font.name = FONT
        rn.font.size = Pt(ln.get('sz', dsz))
        rn.font.bold = ln.get('bold', False)
        rn.font.italic = ln.get('italic', False)
        c = ln.get('c', dc)
        if c: rn.font.color.rgb = c

def rings(sl):
    cx, cy = 11.4, 0.41
    for r in [0.5, 1.0, 1.55, 2.1, 2.65, 3.2, 3.75]:
        O(sl, cx - r, cy - r, r * 2, r * 2, lc=RED, lw=2.8)

def HDR(sl, title):
    R(sl, 0, 0, 13.33, 0.82, fill=NAVY)
    rings(sl)
    T(sl, title, 0.38, 0.12, 9.0, 0.6, sz=22, bold=True, color=WHITE)

def SEC(sl, text, x=0.38, y=0.9, sz=13):
    T(sl, text, x, y, 12.6, 0.36, sz=sz, bold=True, color=BLUE)

def LOGO(sl):
    O(sl, 0.30, 7.1, 0.30, 0.30, lc=BLUE, lw=2.0)
    O(sl, 0.35, 7.15, 0.20, 0.20, lc=BLUE, lw=2.0)
    O(sl, 0.40, 7.20, 0.10, 0.10, lc=BLUE, lw=2.0)
    T(sl, "Московская верфь", 0.68, 7.1, 1.8, 0.34, sz=9, color=NAVY)

def WH(sl):
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = WHITE

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BL = prs.slide_layouts[6]

# ═══════════════════════════════════════
# СЛАЙД 1 — КОНКУРЕНТЫ
# ═══════════════════════════════════════
s1 = prs.slides.add_slide(BL); WH(s1)
HDR(s1, "Очевидные и неочевидные конкуренты")
SEC(s1, "Очевидные конкуренты", y=0.9)

def big_card(sl, x, y, w, h, title, tag, tag_color, bullets, price):
    R(sl, x, y, w, h, fill=WHITE, lc=C(0xCC, 0xD5, 0xE5), lw=0.8)
    R(sl, x, y, 0.065, h, fill=BLUE)
    T(sl, title, x+0.14, y+0.1, w-1.6, 0.34, sz=11.5, bold=True, color=NAVY)
    tw = max(1.1, len(tag)*0.075 + 0.22)
    R(sl, x+w-tw-0.1, y+0.1, tw, 0.24, fill=tag_color)
    T(sl, tag, x+w-tw-0.1, y+0.11, tw, 0.22,
      sz=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    T(sl, '\n'.join(f'— {b}' for b in bullets),
      x+0.14, y+0.5, w-0.24, h-0.68, sz=10.5, color=MDM)
    T(sl, price, x+0.14, y+h-0.3, w-0.24, 0.26,
      sz=10.5, bold=True, color=BLUE)

big_card(s1, 0.38, 1.2, 6.1, 1.85,
    "Astoria Grande — Miray Cruises / СК «Аквамир»", "Международный", NAVY,
    ["Черноморский круизный продукт из Сочи",
     "Маршрут: Сочи → Турция / Греция / Египет",
     "Навигация: март–январь (2023–2026)",
     "2025: 25 тыс. пасс. в март–август (~15% ↑)"],
    "~80–160 тыс. ₽ / чел.")

big_card(s1, 6.74, 1.2, 6.2, 1.85,
    "Речные круизы по России — «Водоходъ», «Мостурфлот» и др.",
    "Речной", C(0x1D, 0x7A, 0x5A),
    ["Активное расширение рынка в 2025–2026 гг.",
     "Новые маршруты: Новосибирск, Томск, Сургут, Салехард",
     "Развитие круизов на Соловки, Енисей, Байкал"],
    "~60–150 тыс. ₽ / чел.")

SEC(s1, "Неочевидные конкуренты", y=3.14)

small = [
    ("Mriya Resort & SPA (Крым)",
     ["Загрузка ~90%", "2025: «Код наслаждения», «Территория сна»",
      "23 000 участников"], "~25–45 тыс. ₽ / ночь"),
    ("Красная Поляна (Сочи)",
     ["370 тыс. туристов (лето 2025, ~+5%)",
      "600 тыс. с проживанием, 270 тыс. — на 1 день",
      "Бронирование: 2–4 недели"], "~0–25 тыс. ₽ / ночь"),
    ("Геленджик / КМВ / Калининград",
     ["Геленджик: 4,5 млн туристов (2025)",
      "КМВ: 250 тыс. (лето), Кисловодск — 141 тыс./год",
      "Море, горы, санатории"], "~5–35 тыс. ₽ / ночь"),
    ("Карелия",
     ["1,4 млн туристов (2025, ~+40%)",
      "Ruskeala Symphony, «В сердце Карелии» (45 тыс.)",
      "Природа, рыбалка, дети"], "~20–60 тыс. ₽ / тур"),
    ("Винные туры («Виноградный экспресс»)",
     ["Крым и Краснодарский край",
      "Посещение виноделен (Alta Valley, Fanagoria)",
      "Формат: лекции + дегустации"], "~15–50 тыс. ₽ / тур"),
    ("Межрегиональные туристические маршруты",
     ["«Великий Русский Северный Путь» — 15 регионов",
      "«Сибирский экспресс» и другие маршруты",
      "Турпоездки 5–10 дней"], "поезд / авто / экскурсии"),
]

CW, CH = 4.06, 1.54
for i, (t, b, p) in enumerate(small):
    cx = 0.38 + (i % 3) * (CW + 0.17)
    cy = 3.5  + (i // 3) * (CH + 0.14)
    R(s1, cx, cy, CW, CH, fill=GRY)
    R(s1, cx, cy, CW, 0.04, fill=BLUE)
    T(s1, t, cx+0.12, cy+0.1, CW-0.22, 0.3, sz=10.5, bold=True, color=NAVY)
    T(s1, '\n'.join(f'· {x}' for x in b),
      cx+0.12, cy+0.44, CW-0.22, CH-0.64, sz=9.5, color=MDM)
    T(s1, p, cx+0.12, cy+CH-0.3, CW-0.22, 0.26, sz=9.5, bold=True, color=BLUE)

LOGO(s1)

# ═══════════════════════════════════════
# СЛАЙД 2 — ТЕПЛОХОДЫ
# ═══════════════════════════════════════
s2 = prs.slides.add_slide(BL); WH(s2)
HDR(s2, "Новые круизные теплоходы 2026–2027")

R(s2, 7.98, 0.82, 5.35, 6.56, fill=FBKG)
R(s2, 7.97, 0.82, 0.02, 6.56, fill=TLINE)

T(s2, "Ключевые новинки рынка", 0.38, 0.92, 7.4, 0.3,
  sz=12, bold=True, color=BLUE)
T(s2, "Проекты в разработке", 8.1, 0.92, 5.0, 0.3,
  sz=12, bold=True, color=BLUE)

ships = [
    ("«Николай Жарков» — «Водоходъ»",
     ["4 палубы", "176 пасс. (люкс)", "запуск 2026"],
     ["Проект «Карелия» · маршруты: реки и моря России + международные круизы",
      "Заход на Беломорско-Балтийский канал → Соловки",
      "«Михаил Хмаринский», «Расул Гамзатов» — рейсы с 2027"], 1.06),
    ("«Байкал» — «Водоходъ»",
     ["3 палубы", "70/152 пасс.", "Электро", "навигация 2027"],
     ["Проект ТФРП 700 · спуск на воду: 2026",
      "Самый экологичный теплоход в РФ · уровень 5 звезд отель на воде"], 0.82),
    ("Новое судно на Байкале — ГК «Истлэнд»",
     ["4 палубы", "144 каюты / 200 пасс.", "запуск 2027"],
     ["Проект RPV 6714 · станет крупнейшим теплоходом на Байкале"], 0.68),
    ("Маршруты с ограниченными зонами",
     ["310 пасс.", "ориентир 2026"],
     ["Проект PV300VD · заполняемость ~90% (2025)",
      "Международные круизы по Каспию"], 0.82),
    ("«Виктор Астафьев» — Пассажиррычтранс",
     ["4 палубы", "245 пасс.", "2026–2027"],
     ["Проект А45-90.2 · Енисей (Красноярск — Дудинка)",
      "Экспедиционные маршруты · выход в Арктику"], 0.82),
]

sy = 1.3
for name, chips, bullets, sh in ships:
    R(s2, 0.38, sy, 7.42, sh, fill=WHITE, lc=TLINE, lw=0.5)
    R(s2, 0.38, sy, 0.065, sh, fill=BLUE)
    T(s2, name, 0.52, sy+0.07, 4.9, 0.3, sz=10.5, bold=True, color=NAVY)
    cx = 5.5
    for chip in chips:
        cw = max(0.72, len(chip) * 0.07 + 0.2)
        R(s2, cx, sy+0.06, cw, 0.2, fill=C(0xE8, 0xEE, 0xF8))
        T(s2, chip, cx+0.03, sy+0.07, cw-0.04, 0.18,
          sz=7.5, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
        cx += cw + 0.05
    T(s2, '\n'.join(f'· {b}' for b in bullets),
      0.52, sy+0.4, 7.16, sh-0.46, sz=10, color=MDM)
    sy += sh + 0.09

devs = [
    ("Проект «Карелия» — «Доннатурфлот» + ОСК", "3", "судна",
     ["5 палуб", "Круизные лайнеры"],
     ["Соглашение подписано на выставке «Нева»",
      "Расширение линейки пятипалубных круизных лайнеров"]),
    ("Проект «Сириус» — «Созвездие»", "6", "судов",
     ["Волго-Дон Макс", "279 пасс.", "85% с балконами"],
     ["Регионы: Финский залив, Азовское и Черное моря",
      "Панорамные шлюзы в носовой части"]),
]

dy = 1.3
for dname, cnt, unit, dchips, dbullets in devs:
    ph = 2.1
    R(s2, 8.12, dy, 4.96, ph, fill=WHITE, lc=TLINE, lw=0.5)
    R(s2, 8.12, dy, 4.96, 0.042, fill=BLUE)
    T(s2, dname, 8.24, dy+0.1, 3.3, 0.38, sz=10.5, bold=True, color=NAVY)
    T(s2, cnt, 11.6, dy+0.04, 1.3, 0.44,
      sz=26, bold=True, color=BLUE, align=PP_ALIGN.RIGHT)
    T(s2, unit, 11.6, dy+0.48, 1.3, 0.22, sz=9, color=LGT, align=PP_ALIGN.RIGHT)
    cx = 8.24
    for chip in dchips:
        cw = max(0.82, len(chip) * 0.072 + 0.22)
        R(s2, cx, dy+0.62, cw, 0.2, fill=C(0xE8, 0xEE, 0xF8))
        T(s2, chip, cx+0.03, dy+0.63, cw-0.04, 0.18,
          sz=7.5, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
        cx += cw + 0.06
    T(s2, '\n'.join(f'· {b}' for b in dbullets),
      8.24, dy+0.9, 4.7, 1.05, sz=10.5, color=MDM)
    dy += ph + 0.15

LOGO(s2)

# ═══════════════════════════════════════
# СЛАЙД 3 — ПОРТРЕТ КЛИЕНТА
# ═══════════════════════════════════════
s3 = prs.slides.add_slide(BL); WH(s3)
HDR(s3, "Портрет клиента речных круизов")
SEC(s3, "Что хочет клиент: сегментация, мотивы", y=0.9, sz=15)

T(s3, "Возрастная структура клиента*", 0.38, 1.42, 3.0, 0.26, sz=11, color=MDM)

age_data = [
    (NAVY,                 "50+ лет",   "65%"),
    (BLUE,                 "40–49 лет", "15%"),
    (C(0x8F, 0xAD, 0xD4), "30–39 лет", "12%"),
    (C(0xC5, 0xD8, 0xF0), "25–30 лет",  "8%"),
]
ay = 1.76
for ac, label, pct in age_data:
    R(s3, 0.38, ay, 0.22, 0.22, fill=ac)
    T(s3, label, 0.66, ay+0.01, 1.5, 0.22, sz=10.5, color=DRK)
    T(s3, pct, 2.0, ay+0.01, 0.6, 0.22,
      sz=10.5, bold=True, color=NAVY, align=PP_ALIGN.RIGHT)
    ay += 0.3

T(s3, "*На основе данных круизного оператора Водоходъ",
  0.38, 2.9, 2.85, 0.34, sz=8.5, italic=True, color=LGT)
T(s3, "Основной профиль клиента",
  0.38, 3.32, 2.85, 0.28, sz=12, bold=True, color=ORANGE)
T(s3, "Пары и семьи 40–65+,\nзаинтересованные в комфортном,\n"
      "познавательном и организованном отдыхе",
  0.38, 3.64, 2.85, 0.64, sz=10.5, color=DRK)

T(s3, "JTBD — Jobs To Be Done (мотивации клиентов)",
  3.54, 1.42, 9.6, 0.26, sz=10.5, bold=True, color=BLUE)

R(s3, 3.54, 1.72, 9.72, 0.32, fill=C(0xF0, 0xF5, 0xFF))
for cx, cw, ch in [(3.54, 3.12, "Хочу"), (6.66, 3.14, "Чтобы"),
                    (9.80, 3.46, "Поэтому выбираю")]:
    T(s3, ch, cx+0.1, 1.74, cw, 0.28, sz=11, bold=True, color=BLUE)
R(s3, 3.54, 2.04, 9.72, 0.04, fill=BLUE)

jtbd_rows = [
    ("Отдохнуть без сложной логистики",
     "Не заниматься организацией поездки", "Речной круиз"),
    ("Увидеть несколько городов",
     "Совместить отдых и экскурсии", "Экскурсионный круиз"),
    ("Спокойно отдохнуть",
     "Восстановиться и сменить темп", "Спокойный речной круиз"),
    ("Поехать всей семьёй",
     "Всем поколениям было интересно", "Семейный круиз"),
    ("Получить высокий уровень сервиса",
     "Чувствовать статус и комфорт", "Премиальный круиз"),
]
ry = 2.08; rh = 0.3
for i, (w, b, c) in enumerate(jtbd_rows):
    if i % 2 == 0: R(s3, 3.54, ry, 9.72, rh, fill=TALT)
    R(s3, 3.54, ry+rh-0.02, 9.72, 0.02, fill=TLINE)
    T(s3, w, 3.64, ry+0.06, 2.96, rh-0.1, sz=10, color=DRK)
    T(s3, b, 6.76, ry+0.06, 3.0,  rh-0.1, sz=10, color=DRK)
    T(s3, c, 9.90, ry+0.06, 3.2,  rh-0.1, sz=10, color=DRK)
    ry += rh

personas = [
    ("Комфорт без\nлогистики", "Пара: 50-56 лет\nДоход: выше среднего",
     ["Отдых все включено", "Понятная логистика",
      "Удобные каюты", "Организованные экскурсии"],
     "комфортный отдых без необходимости\nсамостоятельно планировать маршрут",
     ["старые суда", "слабый сервис", "медицина на борту"]),
    ("Культурный\nотдых", "Женщина: 60 лет\nДоход: средний",
     ["Познавательный отдых", "Посещение исторических городов",
      "Экскурсионная программа"],
     "возможность посетить несколько городов\nс культурной программой",
     ["короткие остановки", "слабая экскурсионная программа"]),
    ("Тихий\nотдых", "Женщина: 48 лет\nДоход: выше среднего",
     ["Спокойная атмосфера", "Природные маршруты",
      "Медленный темп путешествия"],
     "спокойный отдых и\nвозможность перезагрузки",
     ["шумные развлечения", "ограниченные возможности для спокойного отдыха"]),
    ("Семейный\nотдых", "Семья: 35-65 лет, дети\nДоход: средний",
     ["Отдых для нескольких поколений", "Безопасность",
      "Различные развлечения"],
     "возможность путешествовать с семьёй",
     ["маленькие каюты", "ограниченные развлечения",
      "отсутствие детских программ"]),
    ("Статусный\nотдых", "Мужчина: 45 лет\nДоход: высокий",
     ["Приватность", "Высокий уровень сервиса",
      "Хорошая кухня", "Новые суда"],
     "спокойный отдых и\nвозможность перезагрузки",
     ["старые суда", "шумные развлечения", "отсутствие хорошей кухни"]),
]

PW, PH = 2.44, 3.54
for i, (name, demo, imp, driver, barriers) in enumerate(personas):
    px = 0.38 + i * (PW + 0.1)
    py = 3.82
    R(s3, px, py, PW, PH, fill=GRY)
    R(s3, px, py, PW, 0.04, fill=BLUE)
    T(s3, name, px+0.12, py+0.1, PW-0.22, 0.42, sz=10.5, color=BLUE)
    cy = py + 0.58
    T(s3, demo, px+0.12, cy, PW-0.22, 0.4, sz=10, color=DRK); cy += 0.46
    T(s3, "Что важно:", px+0.12, cy, PW-0.22, 0.2, sz=10, bold=True, color=DRK)
    cy += 0.22
    T(s3, '\n'.join(imp), px+0.12, cy, PW-0.22, 0.74, sz=10, color=DRK)
    cy += 0.78
    T(s3, "Драйвер покупки:", px+0.12, cy, PW-0.22, 0.2, sz=10, bold=True, color=DRK)
    cy += 0.22
    T(s3, driver, px+0.12, cy, PW-0.22, 0.52, sz=10, color=DRK); cy += 0.56
    T(s3, "Барьеры:", px+0.12, cy, PW-0.22, 0.2, sz=10, bold=True, color=DRK)
    cy += 0.22
    T(s3, '\n'.join(barriers), px+0.12, cy, PW-0.22, 0.5, sz=10, color=DRK)

LOGO(s3)

prs.save("moskverf_template.pptx")
print("Готово! Файл: moskverf_template.pptx")
print(f"Слайдов: {len(prs.slides)}")
